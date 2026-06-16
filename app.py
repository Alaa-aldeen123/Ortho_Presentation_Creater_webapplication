import os
import sys
import io
import gc
import cv2
import numpy as np
import streamlit as st
from PIL import Image
from rembg import remove
from pptx import Presentation
from pptx.util import Inches
from ultralytics import YOLO
import torch

# Set PyTorch to use CPU to avoid memory issues in Streamlit Cloud
torch.set_num_threads(1)
if torch.cuda.is_available():
    torch.cuda.set_per_process_memory_fraction(0.1)

st.set_page_config(page_title="Orthodontic Case Presentation", page_icon="🦷", layout="wide")

# Custom CSS for Background Color, File Uploader modifications, and Buttons
st.markdown("""
    <style>
    /* Hide the "Limit 200MB per file" text natively injected by Streamlit */
    [data-testid="stFileUploadDropzone"] small {
        display: none !important;
    }
    /* Styling for the Process button */
    div.stButton > button:first-child {
        background-color: #007BFF;
        color: white;
        border-radius: 8px;
        padding: 10px 24px;
        font-size: 18px;
        font-weight: 600;
        border: none;
        transition: background-color 0.3s ease;
    }
    div.stButton > button:first-child:hover {
        background-color: #0056b3;
        border-color: #0056b3;
    }
    /* Styling for the Download button */
    div.stDownloadButton > button:first-child {
        background-color: #28A745;
        color: white;
        border-radius: 8px;
        padding: 10px 24px;
        font-size: 18px;
        font-weight: 600;
        border: none;
        width: 100%;
        transition: background-color 0.3s ease;
    }
    div.stDownloadButton > button:first-child:hover {
        background-color: #218838;
        border-color: #218838;
    }
    </style>
""", unsafe_allow_html=True)

# -------------------------------------------------------------------
# Helper: resource_path()
# Returns the absolute path to a resource whether running as a script
# or as a bundled executable via PyInstaller.
def resource_path(relative_path):
    if getattr(sys, 'frozen', False):  # running in a bundle
        base_path = sys._MEIPASS
    else:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)

# -------------------------------------------------------------------
# ✅ FIX #1 & #2: Load heavy models ONCE with @st.cache_resource.
# This prevents loading YOLO / face-detector dozens of times and
# keeps memory within Streamlit Cloud's 1 GB limit.

@st.cache_resource
def load_yolo_model(model_path):
    """Loads the YOLO model once and reuses it across all calls."""
    return YOLO(model_path)

@st.cache_resource
def load_face_detector(pb_path, pbtxt_path):
    """Loads the OpenCV DNN face detector once."""
    return cv2.dnn.readNetFromTensorflow(pb_path, pbtxt_path)

# -------------------------------------------------------------------
# Use resource_path to get the correct paths for the resource files
pptx_resource = resource_path("main.pptx")
model_resource = resource_path("best.pt")
face_detector_pb = resource_path("opencv_face_detector_uint8.pb")
face_detector_pbtxt = resource_path("opencv_face_detector.pbtxt")

# ✅ FIX #3: Verify all resource files exist at startup with clear errors.
_missing = []
for _label, _path in [("PowerPoint template", pptx_resource),
                      ("YOLO model", model_resource),
                      ("Face detector model", face_detector_pb),
                      ("Face detector config", face_detector_pbtxt)]:
    if not os.path.isfile(_path):
        _missing.append(f"• {_label}: `{_path}`")

if _missing:
    st.error("⚠️ **Missing resource files.** The following files were not found in the repository:\n\n"
             + "\n".join(_missing)
             + "\n\nMake sure they are committed to your GitHub repo **root** (next to `app.py`).")
    st.stop()

# Load the PPTX template into memory using the bundled resource.
try:
    with open(pptx_resource, "rb") as file:
        pptx_data = file.read()
except FileNotFoundError:
    st.error(f"⚠️ Template file not found: {pptx_resource}")
    st.stop()

# Global variables that the processing functions rely on.
patient_name = ""
folder_path = ""

# -------------------------------------------------------------------
# Processing Functions

def create_new_presentation():
    """Create a new presentation based on the stored template data."""
    pp_name = os.path.splitext(patient_name.replace(' ', '_'))[0] + '.pptx'
    pptx_path_local = os.path.join(folder_path, pp_name)
    os.makedirs(folder_path, exist_ok=True)
    if not os.path.exists(pptx_path_local):
        pptx_stream = io.BytesIO(pptx_data)
        new_presentation = Presentation(pptx_stream)
        new_presentation.save(pptx_path_local)

def copy_and_rename_convert_images():
    """
    Copies and renames only the images in the uploaded_image_paths dictionary,
    converting them to JPG.
    """
    for var_name, image_path in uploaded_image_paths.items():
        if os.path.isfile(image_path) and image_path.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff')):
            try:
                copied_image_name = f"{var_name}.jpg"
                copied_image_path = os.path.join(folder_path, copied_image_name)
                image = Image.open(image_path).convert("RGB")
                image.save(copied_image_path, "JPEG")
            except Exception as e:
                st.write(f"Error processing {image_path}: {e}")

def remove_background(image_name):
    """Removes the background from an image and overwrites the file."""
    image_path = os.path.join(folder_path, image_name)
    if os.path.isfile(image_path) and image_name.lower().endswith(('.png', '.jpg', '.jpeg')):
        try:
            input_image = Image.open(image_path)
            output_image = remove(input_image)
            white_background = Image.new("RGB", output_image.size, (255, 255, 255))
            # Guard against images without an alpha channel.
            if output_image.mode == 'RGBA':
                white_background.paste(output_image, mask=output_image.split()[3])
            else:
                white_background.paste(output_image)
            white_background.save(image_path)
        except Exception as e:
            st.write(f"Error processing {image_name}: {e}")

def crop_personal(image_name, left_padding=0.2, right_padding=0.2, above_padding=0.3, bottom_padding=0.1):
    """Crops personal images using face detection."""
    # ✅ FIX #2: Use the cached face detector instead of reloading each call.
    net = load_face_detector(face_detector_pb, face_detector_pbtxt)
    image_path = os.path.join(folder_path, image_name)
    img = cv2.imread(image_path)
    if img is None:
        return image_path
    h, w = img.shape[:2]
    blob = cv2.dnn.blobFromImage(img, scalefactor=1.0, size=(300, 300), mean=(104.0, 177.0, 123.0))
    net.setInput(blob)
    detections = net.forward()
    best_face = None
    for i in range(detections.shape[2]):
        confidence = detections[0, 0, i, 2]
        if confidence > 0.5:
            box = detections[0, 0, i, 3:7] * np.array([w, h, w, h])
            x, y, x_max, y_max = box.astype("int")
            best_face = (x, y, x_max - x, y_max - y)
    if best_face:
        x, y, box_w, box_h = best_face
        top_pad = int(box_h * above_padding)
        bottom_pad = int(box_h * bottom_padding)
        left_pad = int(box_w * left_padding)
        right_pad = int(box_w * right_padding)
        y = max(0, y - top_pad)
        box_h = min(img.shape[0] - y, box_h + top_pad + bottom_pad)
        x = max(0, x - left_pad)
        box_w = min(img.shape[1] - x, box_w + left_pad + right_pad)
        cropped_img = img[y:y + box_h, x:x + box_w]
        cv2.imwrite(image_path, cropped_img)
        return image_path
    else:
        return image_path

def crop_arch(image_name):
    """Crops arch images using a YOLO model."""
    image_path = os.path.join(folder_path, image_name)
    try:
        # ✅ FIX #1: Use the cached YOLO model instead of reloading each call.
        model = load_yolo_model(model_resource)
        image = cv2.imread(image_path)
        if image is None:
            st.write(f"Error: Unable to read image {image_name}")
            return
        results = model(image_path, verbose=False)
        if results and results[0].boxes.xyxy.shape[0] > 0:
            x1, y1, x2, y2 = map(int, results[0].boxes.xyxy[0])
            cropped_object = image[y1:y2, x1:x2]
            cv2.imwrite(image_path, cropped_object)
        else:
            pass
    except Exception as e:
        st.warning(f"Could not process arch image {image_name}: {e}")

def resize_image(image_name, max_width=None, max_height=None, file_prefix=''):
    """Resizes an image to given dimensions and saves a new version with an optional prefix."""
    image_path = os.path.join(folder_path, image_name)
    img = Image.open(image_path)
    orig_width, orig_height = img.size
    dpi = 96
    max_width_px = max_width * dpi if max_width else None
    max_height_px = max_height * dpi if max_height else None
    if max_width_px and max_height_px:
        scale = min(max_width_px / orig_width, max_height_px / orig_height)
    elif max_width_px:
        scale = max_width_px / orig_width
    elif max_height_px:
        scale = max_height_px / orig_height
    else:
        scale = 1
    new_width = int(orig_width * scale)
    new_height = int(orig_height * scale)
    resized_img = img.resize((new_width, new_height), Image.LANCZOS)
    base_name, ext = os.path.splitext(image_name)
    new_file_name = f"{file_prefix}{base_name}{ext}"
    temp_path = os.path.join(folder_path, new_file_name)
    resized_img.save(temp_path)
    return new_file_name

def insert_image(slide_index, image_name, left=None, bottom=None, right=None, top=None):
    """Inserts an image into a slide of the PowerPoint presentation."""
    pptx_path_local = os.path.join(folder_path, os.path.splitext(patient_name.replace(' ', '_'))[0] + '.pptx')
    ppt = Presentation(pptx_path_local)
    slide = ppt.slides[slide_index]
    image_path = os.path.join(folder_path, image_name)
    with Image.open(image_path) as img:
        dpi_info = img.info.get('dpi', (96, 96))
        if isinstance(dpi_info, tuple):
            if len(dpi_info) < 2:
                if len(dpi_info) == 1:
                    dpi = (dpi_info[0], dpi_info[0])
                else:
                    dpi = (96, 96)
            else:
                dpi = dpi_info
        elif isinstance(dpi_info, (int, float)):
            dpi = (dpi_info, dpi_info)
        else:
            dpi = (96, 96)
        width_in_inches = img.width / dpi[0]
        height_in_inches = img.height / dpi[1]
    slide_width = ppt.slide_width
    slide_height = ppt.slide_height
    if right is not None:
        left_val = slide_width - Inches(right) - Inches(width_in_inches)
    elif left is not None:
        left_val = Inches(left)
    else:
        left_val = Inches(0)
    if bottom is not None:
        top_val = slide_height - Inches(bottom) - Inches(height_in_inches)
    elif top is not None:
        top_val = Inches(top)
    else:
        top_val = Inches(0)
    slide.shapes.add_picture(image_path, left_val, top_val,
                             width=Inches(width_in_inches),
                             height=Inches(height_in_inches))
    ppt.save(pptx_path_local)

def run_all_processing():
    """Runs the complete processing pipeline."""
    copy_and_rename_convert_images()
    for img in ['pre_personal_front.jpg', 'pre_personal_smile.jpg', 'pre_personal_oblique.jpg',
                'pre_personal_profile.jpg', 'post_personal_front.jpg', 'post_personal_smile.jpg',
                'post_personal_oblique.jpg', 'post_personal_profile.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            remove_background(img)
    for img in ['pre_personal_front.jpg', 'pre_personal_smile.jpg', 'post_personal_front.jpg',
                'post_personal_smile.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            crop_personal(img)
    for img in ['pre_personal_oblique.jpg', 'post_personal_oblique.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            crop_personal(img, left_padding=0.5)
    for img in ['pre_personal_profile.jpg', 'post_personal_profile.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            crop_personal(img, left_padding=1)
    for img in ['pre_arch_right.jpg', 'pre_arch_front.jpg', 'pre_arch_left.jpg', 'pre_arch_upper.jpg',
                'pre_arch_lower.jpg', 'post_arch_right.jpg', 'post_arch_front.jpg', 'post_arch_left.jpg',
                'post_arch_upper.jpg', 'post_arch_lower.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            crop_arch(img)
    for img in ['pre_personal_front.jpg', 'pre_personal_smile.jpg', 'pre_personal_oblique.jpg',
                'pre_personal_profile.jpg', 'post_personal_front.jpg', 'post_personal_smile.jpg',
                'post_personal_oblique.jpg', 'post_personal_profile.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            resize_image(img, 2.9, 3.9)
            resize_image(img, 2.31, 3.01, 'g_')
    for img in ['pre_arch_right.jpg', 'pre_arch_front.jpg', 'pre_arch_left.jpg',
                'post_arch_right.jpg', 'post_arch_front.jpg', 'post_arch_left.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            resize_image(img, 3.53, 2.17)
            resize_image(img, 2.7, 1.77, 'g_')
    for img in ['pre_arch_upper.jpg', 'pre_arch_lower.jpg', 'post_arch_upper.jpg', 'post_arch_lower.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            resize_image(img, 3.74, 2.76)
            resize_image(img, 2.69, 1.95, 'g_')
    for key in ["pre_cast_right", "pre_cast_front", "pre_cast_left"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 3.53, 2.17)
    for key in ["pre_cast_upper", "pre_cast_lower"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 3.74, 2.76)
    for key in ["post_cast_right", "post_cast_front", "post_cast_left"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 3.53, 2.17)
    for key in ["post_cast_upper", "post_cast_lower"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 3.74, 2.76)
    for key in ["pre-Panoramic X-Ray"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 9.5, 4.85)
    for key in ["pre-Cephalometric X-Ray", "pre-Cephalometric X-Ray Tracing"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 5.97, 4.96)
    for key in ["post-Panoramic X-Ray"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 9.5, 4.85)
    for key in ["post-Cephalometric X-Ray", "post-Cephalometric X-Ray Tracing"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 5.97, 4.96)
    if os.path.exists(os.path.join(folder_path, 'pre_personal_front.jpg')):
        insert_image(2, 'pre_personal_front.jpg', 0.59, 1.44)
    if os.path.exists(os.path.join(folder_path, 'pre_personal_smile.jpg')):
        insert_image(2, 'pre_personal_smile.jpg', 3.67, 1.44)
    if os.path.exists(os.path.join(folder_path, 'pre_personal_oblique.jpg')):
        insert_image(2, 'pre_personal_oblique.jpg', 6.72, 1.44)
    if os.path.exists(os.path.join(folder_path, 'pre_personal_profile.jpg')):
        insert_image(2, 'pre_personal_profile.jpg', 9.78, 1.44)
    if os.path.exists(os.path.join(folder_path, 'post_personal_front.jpg')):
        insert_image(9, 'post_personal_front.jpg', 0.6, 1.44)
    if os.path.exists(os.path.join(folder_path, 'post_personal_smile.jpg')):
        insert_image(9, 'post_personal_smile.jpg', 3.67, 1.44)
    if os.path.exists(os.path.join(folder_path, 'post_personal_oblique.jpg')):
        insert_image(9, 'post_personal_oblique.jpg', 6.72, 1.44)
    if os.path.exists(os.path.join(folder_path, 'post_personal_profile.jpg')):
        insert_image(9, 'post_personal_profile.jpg', 9.79, 1.44)
    if os.path.exists(os.path.join(folder_path, 'pre_arch_right.jpg')):
        insert_image(3, 'pre_arch_right.jpg', 1.09, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'pre_arch_front.jpg')):
        insert_image(3, 'pre_arch_front.jpg', 4.87, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'pre_arch_left.jpg')):
        insert_image(3, 'pre_arch_left.jpg', 8.63, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'pre_arch_upper.jpg')):
        insert_image(3, 'pre_arch_upper.jpg', 2.8, top=4.05)
    if os.path.exists(os.path.join(folder_path, 'pre_arch_lower.jpg')):
        insert_image(3, 'pre_arch_lower.jpg', 6.74, top=4.05)
    if os.path.exists(os.path.join(folder_path, 'post_arch_right.jpg')):
        insert_image(10, 'post_arch_right.jpg', 1.09, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'post_arch_front.jpg')):
        insert_image(10, 'post_arch_front.jpg', 4.86, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'post_arch_left.jpg')):
        insert_image(10, 'post_arch_left.jpg', 8.63, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'post_arch_upper.jpg')):
        insert_image(10, 'post_arch_upper.jpg', 2.79, top=4.05)
    if os.path.exists(os.path.join(folder_path, 'post_arch_lower.jpg')):
        insert_image(10, 'post_arch_lower.jpg', 6.74, top=4.05)
    if os.path.exists(os.path.join(folder_path, 'g_pre_personal_front.jpg')):
        insert_image(4, 'g_pre_personal_front.jpg', 0.9, 4.2)
    if os.path.exists(os.path.join(folder_path, 'g_pre_personal_smile.jpg')):
        insert_image(4, 'g_pre_personal_smile.jpg', 3.98, 4.2)
    if os.path.exists(os.path.join(folder_path, 'g_pre_personal_oblique.jpg')):
        insert_image(4, 'g_pre_personal_oblique.jpg', 7.03, 4.2)
    if os.path.exists(os.path.join(folder_path, 'g_pre_personal_profile.jpg')):
        insert_image(4, 'g_pre_personal_profile.jpg', 10.09, 4.2)
    if os.path.exists(os.path.join(folder_path, 'g_post_personal_front.jpg')):
        insert_image(11, 'g_post_personal_front.jpg', 0.9, 4.2)
    if os.path.exists(os.path.join(folder_path, 'g_post_personal_smile.jpg')):
        insert_image(11, 'g_post_personal_smile.jpg', 3.98, 4.2)
    if os.path.exists(os.path.join(folder_path, 'g_post_personal_oblique.jpg')):
        insert_image(11, 'g_post_personal_oblique.jpg', 7.03, 4.2)
    if os.path.exists(os.path.join(folder_path, 'g_post_personal_profile.jpg')):
        insert_image(11, 'g_post_personal_profile.jpg', 10.09, 4.2)
    if os.path.exists(os.path.join(folder_path, 'g_pre_arch_right.jpg')):
        insert_image(4, 'g_pre_arch_right.jpg', 2.01, top=3.46)
    if os.path.exists(os.path.join(folder_path, 'g_pre_arch_front.jpg')):
        insert_image(4, 'g_pre_arch_front.jpg', 4.93, top=3.45)
    if os.path.exists(os.path.join(folder_path, 'g_pre_arch_left.jpg')):
        insert_image(4, 'g_pre_arch_left.jpg', 7.84, top=3.45)
    if os.path.exists(os.path.join(folder_path, 'g_pre_arch_upper.jpg')):
        insert_image(4, 'g_pre_arch_upper.jpg', 3.42, top=5.35)
    if os.path.exists(os.path.join(folder_path, 'g_pre_arch_lower.jpg')):
        insert_image(4, 'g_pre_arch_lower.jpg', 6.44, top=5.35)
    if os.path.exists(os.path.join(folder_path, 'g_post_arch_right.jpg')):
        insert_image(11, 'g_post_arch_right.jpg', 2.01, top=3.46)
    if os.path.exists(os.path.join(folder_path, 'g_post_arch_front.jpg')):
        insert_image(11, 'g_post_arch_front.jpg', 4.93, top=3.45)
    if os.path.exists(os.path.join(folder_path, 'g_post_arch_left.jpg')):
        insert_image(11, 'g_post_arch_left.jpg', 7.84, top=3.45)
    if os.path.exists(os.path.join(folder_path, 'g_post_arch_upper.jpg')):
        insert_image(11, 'g_post_arch_upper.jpg', 3.42, top=5.35)
    if os.path.exists(os.path.join(folder_path, 'g_post_arch_lower.jpg')):
        insert_image(11, 'g_post_arch_lower.jpg', 6.44, top=5.35)
    if os.path.exists(os.path.join(folder_path, 'pre_cast_right.jpg')):
        insert_image(5, 'pre_cast_right.jpg', 1.09, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'pre_cast_front.jpg')):
        insert_image(5, 'pre_cast_front.jpg', 4.87, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'pre_cast_left.jpg')):
        insert_image(5, 'pre_cast_left.jpg', 8.63, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'pre_cast_upper.jpg')):
        insert_image(5, 'pre_cast_upper.jpg', 2.8, top=4.05)
    if os.path.exists(os.path.join(folder_path, 'pre_cast_lower.jpg')):
        insert_image(5, 'pre_cast_lower.jpg', 6.74, top=4.05)
    if os.path.exists(os.path.join(folder_path, 'post_cast_right.jpg')):
        insert_image(12, 'post_cast_right.jpg', 1.09, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'post_cast_front.jpg')):
        insert_image(12, 'post_cast_front.jpg', 4.87, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'post_cast_left.jpg')):
        insert_image(12, 'post_cast_left.jpg', 8.63, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'post_cast_upper.jpg')):
        insert_image(12, 'post_cast_upper.jpg', 2.8, top=4.05)
    if os.path.exists(os.path.join(folder_path, 'post_cast_lower.jpg')):
        insert_image(12, 'post_cast_lower.jpg', 6.74, top=4.05)
    if os.path.exists(os.path.join(folder_path, 'pre-Panoramic X-Ray.jpg')):
        insert_image(6, 'pre-Panoramic X-Ray.jpg', left=1.92, top=1.33)
    if os.path.exists(os.path.join(folder_path, 'pre-Cephalometric X-Ray.jpg')):
        insert_image(7, 'pre-Cephalometric X-Ray.jpg', left=0.55, top=1.5)
    if os.path.exists(os.path.join(folder_path, 'pre-Cephalometric X-Ray Tracing.jpg')):
        insert_image(7, 'pre-Cephalometric X-Ray Tracing.jpg', left=6.81, top=1.5)
    if os.path.exists(os.path.join(folder_path, 'post-Panoramic X-Ray.jpg')):
        insert_image(13, 'post-Panoramic X-Ray.jpg', left=1.92, top=1.33)
    if os.path.exists(os.path.join(folder_path, 'post-Cephalometric X-Ray.jpg')):
        insert_image(14, 'post-Cephalometric X-Ray.jpg', left=0.55, top=1.5)
    if os.path.exists(os.path.join(folder_path, 'post-Cephalometric X-Ray Tracing.jpg')):
        insert_image(14, 'post-Cephalometric X-Ray Tracing.jpg', left=6.81, top=1.5)

# -------------------------------------------------------------------
# Sidebar: Clinic Logo and Instructions
with st.sidebar:
    # Adding your specific clinic logo
    try:
        st.image("logo_colored_with_words.png", use_container_width=True)
    except:
        st.info("🏥 Orthodontic Case Presentation")
    st.markdown("---")
    st.markdown("### 💡 Instructions")
    st.info(
        "**1.** Enter the patient's name in the main window.\n\n"
        "**2.** Navigate through the **Pre** and **Post** treatment tabs to upload images.\n\n"
        "**3.** Once finished, click **Process Images**.\n\n"
        "**4.** Download your Processed File."
    )
    st.markdown("---")
    st.caption("Created by Alaa-Aldeen Aboarrijal 771558504")

# -------------------------------------------------------------------
# Main Title and Description
st.markdown("<h1 style='text-align: center; color: #2C3E50;'>🦷 Orthodontic Case Presentation</h1>",
            unsafe_allow_html=True)
st.markdown(
    "<p style='text-align: center; color: #7F8C8D; font-size: 16px;'>"
    "This application processes Orthodontic Case images to generate a customized presentation for Orthodontists<br>"
    "<b>Created by Alaa-Aldeen Aboarrijal 771558504</b>"
    "</p>",
    unsafe_allow_html=True
)
st.markdown("---")

# -------------------------------------------------------------------
# Patient Details section moved out of sidebar
st.markdown("### 📋 Patient Details")
patient_name = st.text_input("Enter Patient's Name", placeholder="e.g., John Doe")
st.markdown("<br>", unsafe_allow_html=True)  # Adding a bit of spacing

# -------------------------------------------------------------------
# Tabs for Image Uploads
categories = {
    "Pre-Personal Images": ["pre_personal_front", "pre_personal_smile", "pre_personal_oblique", "pre_personal_profile"],
    "Pre-Arch Images": ["pre_arch_right", "pre_arch_front", "pre_arch_left", "pre_arch_upper", "pre_arch_lower"],
    "Pre-Cast Images": ["pre_cast_right", "pre_cast_front", "pre_cast_left", "pre_cast_upper", "pre_cast_lower"],
    "Pre-X-rays Images": ["pre-Panoramic X-Ray", "pre-Cephalometric X-Ray", "pre-Cephalometric X-Ray Tracing"],
    "Post-Personal Images": ["post_personal_front", "post_personal_smile", "post_personal_oblique",
                             "post_personal_profile"],
    "Post-Arch Images": ["post_arch_right", "post_arch_front", "post_arch_left", "post_arch_upper", "post_arch_lower"],
    "Post-Cast Images": ["post_cast_right", "post_cast_front", "post_cast_left", "post_cast_upper", "post_cast_lower"],
    "Post-X-rays Images": ["post-Panoramic X-Ray", "post-Cephalometric X-Ray", "post-Cephalometric X-Ray Tracing"]
}

uploaded_files = {}
tab1, tab2 = st.tabs(["⏳ Pre-Treatment Images", "✨ Post-Treatment Images"])

for cat, keys in categories.items():
    # Determine the target tab based on the category name
    target_tab = tab1 if "Pre" in cat else tab2
    with target_tab:
        with st.expander(f"📁 {cat}", expanded=False):
            cols = st.columns(2)
            for idx, key in enumerate(keys):
                with cols[idx % 2]:
                    # Make the label slightly prettier
                    display_label = key.replace('_', ' ').replace('-', ' ').title()
                    uploaded_files[key] = st.file_uploader(
                        f"Upload {display_label} Image",
                        key=key
                    )

# -------------------------------------------------------------------
# Define a global dictionary to store image paths from uploaded files.
uploaded_image_paths = {}
st.markdown("---")

# Center the button using columns
col1, col2, col3 = st.columns([1, 2, 1])
with col2:
    process_clicked = st.button("🚀 Process Images", use_container_width=True)

if process_clicked:
    if patient_name.strip() == "":
        st.error("⚠️ Please enter the patient's name in the Patient Details section before processing!")
    else:
        import shutil
        import zipfile

        # --- ORGANIZE OUTPUT FOLDER LOGIC ---
        output_base = "output"
        patient_folder = os.path.join(output_base, patient_name)
        folder_path = patient_folder
        os.makedirs(folder_path, exist_ok=True)

        # Create raw and processed folders
        raw_folder = os.path.join(patient_folder, "images", "raw_images")
        processed_folder = os.path.join(patient_folder, "images", "processed_images")

        key_to_category = {}
        for cat, keys in categories.items():
            folder_name = cat.replace(" ", "_")
            for key in keys:
                key_to_category[key] = folder_name
            os.makedirs(os.path.join(raw_folder, folder_name), exist_ok=True)
            os.makedirs(os.path.join(processed_folder, folder_name), exist_ok=True)

        # Save uploaded files to the new raw folders and map their paths
        for key, file in uploaded_files.items():
            if file is not None:
                cat_folder = key_to_category.get(key, "Uncategorized")
                raw_cat_folder = os.path.join(raw_folder, cat_folder)
                file_path = os.path.join(raw_cat_folder, file.name)
                with open(file_path, "wb") as f:
                    f.write(file.read())
                # Assign the raw paths so copy_and_rename_convert_images() can read them
                uploaded_image_paths[key] = file_path
                globals()[key] = file_path

        with st.status("⚙️ Processing images...", expanded=True) as status:
            progress_bar = st.progress(0)
            # Step 1: Create new presentation
            create_new_presentation()
            progress_bar.progress(20)
            # Step 2: Run full processing pipeline
            try:
                run_all_processing()
                # --- POST-PROCESSING ORGANIZATION ---
                for f in os.listdir(folder_path):
                    if f.lower().endswith(('.jpg', '.jpeg', '.png', '.bmp', '.tiff')):
                        base_name = f
                        if base_name.startswith('g_'):
                            base_name = base_name[2:]
                        base_name = os.path.splitext(base_name)[0]
                        cat_folder_name = key_to_category.get(base_name)
                        if cat_folder_name:
                            dest_folder = os.path.join(processed_folder, cat_folder_name)
                            shutil.move(os.path.join(folder_path, f), os.path.join(dest_folder, f))
                progress_bar.progress(100)
                status.update(label="Processing completed successfully!", state="complete", expanded=False)
            except Exception as e:
                progress_bar.progress(0)
                status.update(label=f"Error during processing: {str(e)}", state="error", expanded=True)
                st.error(f"Processing failed: {str(e)}")

        # Step 3: Zip the output folder for download
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
            for root, dirs, files in os.walk(patient_folder):
                for f in files:
                    full_path = os.path.join(root, f)
                    arcname = os.path.relpath(full_path, patient_folder)
                    zf.write(full_path, arcname=arcname)
        zip_buffer.seek(0)

        # ✅ FIX #5: Free memory / clean up temp files after zipping.
        del patient_folder
        gc.collect()

        # Display the download button prominently
        col_dl1, col_dl2, col_dl3 = st.columns([1, 2, 1])
        with col_dl2:
            st.download_button("📥 Download Processed File", data=zip_buffer,
                               file_name=f"{patient_name}.zip", mime="application/zip")

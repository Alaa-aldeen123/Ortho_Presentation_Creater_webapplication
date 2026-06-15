import os
import sys
import io
import cv2
import numpy as np
import tkinter as tk
from tkinter import filedialog, messagebox
from tkinter import ttk
from PIL import Image
from rembg import remove
from pptx import Presentation
from pptx.util import Inches
from ultralytics import YOLO

# -------------------------------------------------------------------
# Helper: resource_path()
# Returns the absolute path to a resource whether running as a script
# or as a bundled executable via PyInstaller.
def resource_path(relative_path):
    if getattr(sys, 'frozen', False):  # running in a bundle
        base_path = sys._MEIPASS
    else:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)

# -------------------------------------------------------------------
# Use resource_path to get the correct paths for the resource files
pptx_resource = resource_path("main.pptx")
model_resource = resource_path("best.pt")

# Load the PPTX template into memory using the bundled resource.
with open(pptx_resource, "rb") as file:
    pptx_data = file.read()

# Global variables that the processing functions rely on.
patient_name = ""
folder_path = ""

# -------------------------------------------------------------------
# Processing Functions

def create_new_presentation():
    """Create a new presentation based on the stored template data."""
    pp_name = os.path.splitext(patient_name.replace(' ', '_'))[0] + '.pptx'
    pptx_path_local = os.path.join(folder_path, pp_name)
    os.makedirs(folder_path, exist_ok=True)
    if not os.path.exists(pptx_path_local):
        pptx_stream = io.BytesIO(pptx_data)
        new_presentation = Presentation(pptx_stream)
        new_presentation.save(pptx_path_local)
        print(f"Created new presentation: {pp_name}")
    else:
        print(f"Presentation {pp_name} already exists.")

def copy_and_rename_convert_images():
    """
    Finds all image path variables (set via globals by the GUI),
    copies them to the output folder, renames them based on their variable names,
    and converts them to JPG.
    """
    for var_name, image_path in globals().items():
        if isinstance(image_path, str) and os.path.isfile(image_path) and image_path.lower().endswith(
                ('.png', '.jpg', '.jpeg', '.bmp', '.tiff')):
            try:
                copied_image_name = f"{var_name}.jpg"
                copied_image_path = os.path.join(folder_path, copied_image_name)
                image = Image.open(image_path).convert("RGB")
                image.save(copied_image_path, "JPEG")
                print(f"Copied, renamed, and converted: {copied_image_name}")
            except Exception as e:
                print(f"Error processing {image_path}: {e}")

def remove_background(image_name):
    """Removes the background from an image and overwrites the file."""
    image_path = os.path.join(folder_path, image_name)
    if os.path.isfile(image_path) and image_name.lower().endswith(('.png', '.jpg', '.jpeg')):
        try:
            input_image = Image.open(image_path)
            output_image = remove(input_image)
            white_background = Image.new("RGB", output_image.size, (255, 255, 255))
            white_background.paste(output_image, mask=output_image.split()[3])
            white_background.save(image_path)
            print(f"Background removed: {image_name}")
        except Exception as e:
            print(f"Error processing {image_name}: {e}")

def crop_personal(image_name, left_padding=0.2, right_padding=0.2, above_padding=0.3, bottom_padding=0.1):
    """Crops personal images using face detection."""
    #model_path_local = cv2.data.haarcascades.replace("haarcascades", "") + "opencv_face_detector_uint8.pb"
    #config_path = cv2.data.haarcascades.replace("haarcascades", "") + "opencv_face_detector.pbtxt"
    #model_path_local = os.path.join(cv2.data.haarcascades, "opencv_face_detector_uint8.pb")
    #config_path = os.path.join(cv2.data.haarcascades, "opencv_face_detector.pbtxt")
    #face_detector_pb = r".venv\Lib\site-packages\cv2\data\opencv_face_detector_uint8.pb"
    face_detector_pb = r'D:\AI cours by Hasoob\Real_Projects\Ortho_Presentation_Creater_desktop\.venv\Lib\site-packages\cv2\data\opencv_face_detector_uint8.pb'
    #face_detector_pbtxt = r".venv\Lib\site-packages\cv2\data\opencv_face_detector.pbtxt"
    face_detector_pbtxt = r'D:\AI cours by Hasoob\Real_Projects\Ortho_Presentation_Creater_desktop\.venv\Lib\site-packages\cv2\data\opencv_face_detector.pbtxt'

    #net = cv2.dnn.readNetFromTensorflow(model_path_local, config_path)
    net = cv2.dnn.readNetFromTensorflow(face_detector_pb, face_detector_pbtxt)
    image_path = os.path.join(folder_path, image_name)
    img = cv2.imread(image_path)
    if img is None:
        print(f"Error: Unable to read image {image_path}")
        return image_path
    h, w = img.shape[:2]
    blob = cv2.dnn.blobFromImage(img, scalefactor=1.0, size=(300, 300), mean=(104.0, 177.0, 123.0))
    net.setInput(blob)
    detections = net.forward()
    best_face = None
    for i in range(detections.shape[2]):
        confidence = detections[0, 0, i, 2]
        if confidence > 0.5:
            box = detections[0, 0, i, 3:7] * np.array([w, h, w, h])
            x, y, x_max, y_max = box.astype("int")
            best_face = (x, y, x_max - x, y_max - y)
    if best_face:
        x, y, box_w, box_h = best_face
        top_pad = int(box_h * above_padding)
        bottom_pad = int(box_h * bottom_padding)
        left_pad = int(box_w * left_padding)
        right_pad = int(box_w * right_padding)
        y = max(0, y - top_pad)
        box_h = min(img.shape[0] - y, box_h + top_pad + bottom_pad)
        x = max(0, x - left_pad)
        box_w = min(img.shape[1] - x, box_w + left_pad + right_pad)
        cropped_img = img[y:y + box_h, x:x + box_w]
        cv2.imwrite(image_path, cropped_img)
        print(f"Photo cropped: {image_name}")
        return image_path
    else:
        print("No face detected")
        return image_path

def crop_arch(image_name):
    """Crops arch images using a YOLO model."""
    image_path = os.path.join(folder_path, image_name)
    # Use the bundled YOLO model resource.
    model = YOLO(model_resource)
    image = cv2.imread(image_path)
    if image is None:
        print(f"Error: Unable to read image {image_name}")
        return
    results = model(image_path, verbose=False)
    if results and results[0].boxes.xyxy.shape[0] > 0:
        x1, y1, x2, y2 = map(int, results[0].boxes.xyxy[0])
        cropped_object = image[y1:y2, x1:x2]
        cv2.imwrite(image_path, cropped_object)
        print(f"Photo cropped: {image_name}")
    else:
        print(f"No objects detected in {image_name}")

def resize_image(image_name, max_width=None, max_height=None, file_prefix=''):
    """Resizes an image to given dimensions and saves a new version with an optional prefix."""
    image_path = os.path.join(folder_path, image_name)
    img = Image.open(image_path)
    orig_width, orig_height = img.size
    dpi = 96
    max_width_px = max_width * dpi if max_width else None
    max_height_px = max_height * dpi if max_height else None
    if max_width_px and max_height_px:
        scale = min(max_width_px / orig_width, max_height_px / orig_height)
    elif max_width_px:
        scale = max_width_px / orig_width
    elif max_height_px:
        scale = max_height_px / orig_height
    else:
        scale = 1
    new_width = int(orig_width * scale)
    new_height = int(orig_height * scale)
    resized_img = img.resize((new_width, new_height), Image.LANCZOS)
    base_name, ext = os.path.splitext(image_name)
    new_file_name = f"{file_prefix}{base_name}{ext}"
    temp_path = os.path.join(folder_path, new_file_name)
    resized_img.save(temp_path)
    print(f"Resized image: {new_file_name}")
    return new_file_name

def insert_image(slide_index, image_name, left=None, bottom=None, right=None, top=None):
    """Inserts an image into a slide of the PowerPoint presentation."""
    pptx_path_local = os.path.join(folder_path, os.path.splitext(patient_name.replace(' ', '_'))[0] + '.pptx')
    ppt = Presentation(pptx_path_local)
    slide = ppt.slides[slide_index]
    image_path = os.path.join(folder_path, image_name)
    with Image.open(image_path) as img:
        dpi_info = img.info.get('dpi', (96, 96))
        # Robust handling of dpi_info: if it's not a proper 2-tuple, default to (96,96)
        if isinstance(dpi_info, tuple):
            if len(dpi_info) < 2:
                if len(dpi_info) == 1:
                    dpi = (dpi_info[0], dpi_info[0])
                else:
                    dpi = (96, 96)
            else:
                dpi = dpi_info
        elif isinstance(dpi_info, (int, float)):
            dpi = (dpi_info, dpi_info)
        else:
            dpi = (96, 96)
        width_in_inches = img.width / dpi[0]
        height_in_inches = img.height / dpi[1]
    slide_width = ppt.slide_width
    slide_height = ppt.slide_height
    if right is not None:
        left_val = slide_width - Inches(right) - Inches(width_in_inches)
    elif left is not None:
        left_val = Inches(left)
    else:
        left_val = Inches(0)
    if bottom is not None:
        top_val = slide_height - Inches(bottom) - Inches(height_in_inches)
    elif top is not None:
        top_val = Inches(top)
    else:
        top_val = Inches(0)
    slide.shapes.add_picture(image_path, left_val, top_val,
                             width=Inches(width_in_inches),
                             height=Inches(height_in_inches))
    ppt.save(pptx_path_local)
    print(f"Inserted image: {image_name}")


def run_all_processing():
    """Runs the complete processing pipeline."""
    # Step 1: Copy, rename, and convert images
    copy_and_rename_convert_images()

    # Step 2: Remove background from personal images
    for img in ['pre_personal_front.jpg', 'pre_personal_smile.jpg', 'pre_personal_oblique.jpg',
                'pre_personal_profile.jpg', 'post_personal_front.jpg', 'post_personal_smile.jpg',
                'post_personal_oblique.jpg', 'post_personal_profile.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            remove_background(img)

    # Step 3: Crop personal images
    for img in ['pre_personal_front.jpg', 'pre_personal_smile.jpg', 'post_personal_front.jpg',
                'post_personal_smile.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            crop_personal(img)
    for img in ['pre_personal_oblique.jpg', 'post_personal_oblique.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            crop_personal(img, left_padding=0.5)
    for img in ['pre_personal_profile.jpg', 'post_personal_profile.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            crop_personal(img, left_padding=1)

    # Step 4: Crop arch images
    for img in ['pre_arch_right.jpg', 'pre_arch_front.jpg', 'pre_arch_left.jpg', 'pre_arch_upper.jpg',
                'pre_arch_lower.jpg', 'post_arch_right.jpg', 'post_arch_front.jpg', 'post_arch_left.jpg',
                'post_arch_upper.jpg', 'post_arch_lower.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            crop_arch(img)

    # Step 5: Resize images for personal and arch images
    for img in ['pre_personal_front.jpg', 'pre_personal_smile.jpg', 'pre_personal_oblique.jpg',
                'pre_personal_profile.jpg', 'post_personal_front.jpg', 'post_personal_smile.jpg',
                'post_personal_oblique.jpg', 'post_personal_profile.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            resize_image(img, 2.9, 3.9)
            resize_image(img, 2.31, 3.01, 'g_')
    for img in ['pre_arch_right.jpg', 'pre_arch_front.jpg', 'pre_arch_left.jpg',
                'post_arch_right.jpg', 'post_arch_front.jpg', 'post_arch_left.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            resize_image(img, 3.53, 2.17)
            resize_image(img, 2.7, 1.77, 'g_')
    for img in ['pre_arch_upper.jpg', 'pre_arch_lower.jpg', 'post_arch_upper.jpg', 'post_arch_lower.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            resize_image(img, 3.74, 2.76)
            resize_image(img, 2.69, 1.95, 'g_')

    # --- New Processing for Cast and X-ray Images ---
    for key in ["pre_cast_right", "pre_cast_front", "pre_cast_left"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 3.53, 2.17)
    for key in ["pre_cast_upper", "pre_cast_lower"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 3.74, 2.76)
    for key in ["post_cast_right", "post_cast_front", "post_cast_left"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 3.53, 2.17)
    for key in ["post_cast_upper", "post_cast_lower"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 3.74, 2.76)
    for key in ["pre-Panoramic X-Ray"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 9.5, 4.85)
    for key in ["pre-Cephalometric X-Ray", "pre-Cephalometric X-Ray Tracing"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 5.97, 4.96)
    for key in ["post-Panoramic X-Ray"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 9.5, 4.85)
    for key in ["post-Cephalometric X-Ray", "post-Cephalometric X-Ray Tracing"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 5.97, 4.96)

    # Step 6: Insert images into PowerPoint

    # Personal Images
    if os.path.exists(os.path.join(folder_path, 'pre_personal_front.jpg')):
        insert_image(2, 'pre_personal_front.jpg', 0.59, 1.44)
    if os.path.exists(os.path.join(folder_path, 'pre_personal_smile.jpg')):
        insert_image(2, 'pre_personal_smile.jpg', 3.67, 1.44)
    if os.path.exists(os.path.join(folder_path, 'pre_personal_oblique.jpg')):
        insert_image(2, 'pre_personal_oblique.jpg', 6.72, 1.44)
    if os.path.exists(os.path.join(folder_path, 'pre_personal_profile.jpg')):
        insert_image(2, 'pre_personal_profile.jpg', 9.78, 1.44)
    if os.path.exists(os.path.join(folder_path, 'post_personal_front.jpg')):
        insert_image(9, 'post_personal_front.jpg', 0.6, 1.44)
    if os.path.exists(os.path.join(folder_path, 'post_personal_smile.jpg')):
        insert_image(9, 'post_personal_smile.jpg', 3.67, 1.44)
    if os.path.exists(os.path.join(folder_path, 'post_personal_oblique.jpg')):
        insert_image(9, 'post_personal_oblique.jpg', 6.72, 1.44)
    if os.path.exists(os.path.join(folder_path, 'post_personal_profile.jpg')):
        insert_image(9, 'post_personal_profile.jpg', 9.79, 1.44)

    # Arch Images (Pre: slide index 3, Post: slide index 10)
    if os.path.exists(os.path.join(folder_path, 'pre_arch_right.jpg')):
        insert_image(3, 'pre_arch_right.jpg', 1.09, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'pre_arch_front.jpg')):
        insert_image(3, 'pre_arch_front.jpg', 4.87, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'pre_arch_left.jpg')):
        insert_image(3, 'pre_arch_left.jpg', 8.63, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'pre_arch_upper.jpg')):
        insert_image(3, 'pre_arch_upper.jpg', 2.8, top=4.05)
    if os.path.exists(os.path.join(folder_path, 'pre_arch_lower.jpg')):
        insert_image(3, 'pre_arch_lower.jpg', 6.74, top=4.05)
    if os.path.exists(os.path.join(folder_path, 'post_arch_right.jpg')):
        insert_image(10, 'post_arch_right.jpg', 1.09, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'post_arch_front.jpg')):
        insert_image(10, 'post_arch_front.jpg', 4.86, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'post_arch_left.jpg')):
        insert_image(10, 'post_arch_left.jpg', 8.63, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'post_arch_upper.jpg')):
        insert_image(10, 'post_arch_upper.jpg', 2.79, top=4.05)
    if os.path.exists(os.path.join(folder_path, 'post_arch_lower.jpg')):
        insert_image(10, 'post_arch_lower.jpg', 6.74, top=4.05)

    # Pre-Cast Images (insert into slide index 5)
    if os.path.exists(os.path.join(folder_path, 'pre_cast_right.jpg')):
        insert_image(5, 'pre_cast_right.jpg', 1.09, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'pre_cast_front.jpg')):
        insert_image(5, 'pre_cast_front.jpg', 4.87, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'pre_cast_left.jpg')):
        insert_image(5, 'pre_cast_left.jpg', 8.63, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'pre_cast_upper.jpg')):
        insert_image(5, 'pre_cast_upper.jpg', 2.8, top=4.05)
    if os.path.exists(os.path.join(folder_path, 'pre_cast_lower.jpg')):
        insert_image(5, 'pre_cast_lower.jpg', 6.74, top=4.05)

    # Post-Cast Images (insert into slide index 12)
    if os.path.exists(os.path.join(folder_path, 'post_cast_right.jpg')):
        insert_image(12, 'post_cast_right.jpg', 1.09, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'post_cast_front.jpg')):
        insert_image(12, 'post_cast_front.jpg', 4.87, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'post_cast_left.jpg')):
        insert_image(12, 'post_cast_left.jpg', 8.63, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'post_cast_upper.jpg')):
        insert_image(12, 'post_cast_upper.jpg', 2.8, top=4.05)
    if os.path.exists(os.path.join(folder_path, 'post_cast_lower.jpg')):
        insert_image(12, 'post_cast_lower.jpg', 6.74, top=4.05)

    # Pre-X-rays Images (only resized; insertion based on type)
    if os.path.exists(os.path.join(folder_path, 'pre-Panoramic X-Ray.jpg')):
        insert_image(6, 'pre-Panoramic X-Ray.jpg', left=1.92, top=1.33)
    if os.path.exists(os.path.join(folder_path, 'pre-Cephalometric X-Ray.jpg')):
        insert_image(7, 'pre-Cephalometric X-Ray.jpg', left=0.55, top=1.5)
    if os.path.exists(os.path.join(folder_path, 'pre-Cephalometric X-Ray Tracing.jpg')):
        insert_image(7, 'pre-Cephalometric X-Ray Tracing.jpg', left=6.81, top=1.5)

    # Post-X-rays Images
    if os.path.exists(os.path.join(folder_path, 'post-Panoramic X-Ray.jpg')):
        insert_image(13, 'post-Panoramic X-Ray.jpg', left=1.92, top=1.33)
    if os.path.exists(os.path.join(folder_path, 'post-Cephalometric X-Ray.jpg')):
        insert_image(14, 'post-Cephalometric X-Ray.jpg', left=0.55, top=1.5)
    if os.path.exists(os.path.join(folder_path, 'post-Cephalometric X-Ray Tracing.jpg')):
        insert_image(14, 'post-Cephalometric X-Ray Tracing.jpg', left=6.81, top=1.5)

# -------------------------------------------------------------------
# Improved UI/UX Application

class DentalImageProcessorApp(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("Dental Image Processor")
        self.geometry("800x600")
        self.configure(bg="#f0f0f0")
        self.file_labels = {}
        # Dictionary to store image file paths from the GUI
        self.image_paths = {
            "pre_personal_front": None, "pre_personal_smile": None, "pre_personal_oblique": None, "pre_personal_profile": None,
            "pre_arch_right": None, "pre_arch_front": None, "pre_arch_left": None, "pre_arch_upper": None, "pre_arch_lower": None,
            "pre_cast_right": None, "pre_cast_front": None, "pre_cast_left": None, "pre_cast_upper": None, "pre_cast_lower": None,
            "pre-Panoramic X-Ray": None, "pre-Cephalometric X-Ray": None, "pre-Cephalometric X-Ray Tracing": None,
            "post_personal_front": None, "post_personal_smile": None, "post_personal_oblique": None, "post_personal_profile": None,
            "post_arch_right": None, "post_arch_front": None, "post_arch_left": None, "post_arch_upper": None, "post_arch_lower": None,
            "post_cast_right": None, "post_cast_front": None, "post_cast_left": None, "post_cast_upper": None, "post_cast_lower": None,
            "post-Panoramic X-Ray": None, "post-Cephalometric X-Ray": None, "post-Cephalometric X-Ray Tracing": None
        }
        header_frame = ttk.Frame(self, padding="10")
        header_frame.pack(side="top", fill="x")
        ttk.Label(header_frame, text="Patient's Name:").grid(row=0, column=0, sticky="w", padx=5, pady=5)
        self.patient_name_entry = ttk.Entry(header_frame, width=40)
        self.patient_name_entry.grid(row=0, column=1, sticky="w", padx=5, pady=5)
        ttk.Button(header_frame, text="Select Output Folder", command=self.select_output_folder).grid(row=1, column=0, padx=5, pady=5)
        self.output_folder_label = ttk.Label(header_frame, text="No folder selected", foreground="gray")
        self.output_folder_label.grid(row=1, column=1, sticky="w", padx=5, pady=5)
        self.notebook = ttk.Notebook(self)
        self.notebook.pack(fill="both", expand=True, padx=10, pady=10)
        self.create_tab("Pre-Personal Images",
                        ["pre_personal_front", "pre_personal_smile", "pre_personal_oblique", "pre_personal_profile"])
        self.create_tab("Pre-Arch Images",
                        ["pre_arch_right", "pre_arch_front", "pre_arch_left", "pre_arch_upper", "pre_arch_lower"])
        self.create_tab("Pre-Cast Images",
                        ["pre_cast_right", "pre_cast_front", "pre_cast_left", "pre_cast_upper", "pre_cast_lower"])
        self.create_tab("Pre-X-rays Images",
                        ["pre-Panoramic X-Ray", "pre-Cephalometric X-Ray", "pre-Cephalometric X-Ray Tracing"])
        self.create_tab("Post-Personal Images",
                        ["post_personal_front", "post_personal_smile", "post_personal_oblique", "post_personal_profile"])
        self.create_tab("Post-Arch Images",
                        ["post_arch_right", "post_arch_front", "post_arch_left", "post_arch_upper", "post_arch_lower"])
        self.create_tab("Post-Cast Images",
                        ["post_cast_right", "post_cast_front", "post_cast_left", "post_cast_upper", "post_cast_lower"])
        self.create_tab("Post-X-rays Images",
                        ["post-Panoramic X-Ray", "post-Cephalometric X-Ray", "post-Cephalometric X-Ray Tracing"])
        action_frame = ttk.Frame(self, padding="10")
        action_frame.pack(side="bottom", fill="x")
        self.process_button = ttk.Button(action_frame, text="Process Images", command=self.process_images)
        self.process_button.pack(pady=5)
        self.status_label = ttk.Label(action_frame, text="", foreground="green")
        self.status_label.pack(pady=5)

    def create_tab(self, tab_name, keys):
        frame = ttk.Frame(self.notebook, padding="10")
        self.notebook.add(frame, text=tab_name)
        for i, key in enumerate(keys):
            ttk.Label(frame, text=key.replace("_", " ").capitalize() + ":").grid(row=i, column=0, sticky="w", padx=5, pady=5)
            btn = ttk.Button(frame, text="Select File", command=lambda k=key: self.select_file(k))
            btn.grid(row=i, column=1, padx=5, pady=5)
            label = ttk.Label(frame, text="No file selected", foreground="gray")
            label.grid(row=i, column=2, sticky="w", padx=5, pady=5)
            self.file_labels[key] = label
            self.image_paths[key] = None

    def select_file(self, key):
        file_path = filedialog.askopenfilename(title="Select an Image",
                                               filetypes=[("Image Files", "*.png;*.jpg;*.jpeg;*.bmp;*.tiff")])
        if file_path:
            self.image_paths[key] = file_path
            self.file_labels[key].config(text=os.path.basename(file_path))
        else:
            self.image_paths[key] = None
            self.file_labels[key].config(text="No file selected")

    def select_output_folder(self):
        folder = filedialog.askdirectory(title="Select Output Folder")
        if folder:
            self.output_folder_label.config(text=folder)
        else:
            self.output_folder_label.config(text="No folder selected")

    def process_images(self):
        global patient_name, folder_path
        patient_name = self.patient_name_entry.get().strip()
        output_folder = self.output_folder_label.cget("text")
        if not patient_name:
            messagebox.showerror("Error", "Please enter the patient's name!")
            return
        if output_folder == "No folder selected":
            messagebox.showerror("Error", "Please select an output folder!")
            return
        folder_path = os.path.join(output_folder, patient_name)
        os.makedirs(folder_path, exist_ok=True)
        for key, path in self.image_paths.items():
            if path:
                globals()[key] = path
        self.status_label.config(text="Processing images...", foreground="blue")
        self.update()
        create_new_presentation()
        run_all_processing()
        self.status_label.config(text="Processing complete!", foreground="green")
        messagebox.showinfo("Success", "Processing complete! Check the output folder for the presentation and images.")

if __name__ == "__main__":
    app = DentalImageProcessorApp()
    app.mainloop()
